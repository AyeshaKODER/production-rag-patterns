"""
Multi-format text extraction for ingestion.

Each function takes a filepath and returns plain text — the output feeds
into RAGPipeline's existing chunking logic unchanged, so adding a new
format here never touches chunking, embedding, or retrieval code.
"""

import fitz  # PyMuPDF
from docx import Document as DocxDocument
from pptx import Presentation
from bs4 import BeautifulSoup


def extract_pdf_text(filepath: str) -> str:
    doc = fitz.open(filepath)
    text_parts = [page.get_text() for page in doc]
    doc.close()
    return "\n\n".join(text_parts)


def extract_docx_text(filepath: str) -> str:
    doc = DocxDocument(filepath)
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    return "\n\n".join(paragraphs)


def extract_pptx_text(filepath: str) -> str:
    prs = Presentation(filepath)
    slide_texts = []
    for i, slide in enumerate(prs.slides, start=1):
        texts = [
            shape.text for shape in slide.shapes
            if shape.has_text_frame and shape.text.strip()
        ]
        if texts:
            slide_texts.append(f"## Slide {i}\n" + "\n".join(texts))
    return "\n\n".join(slide_texts)


def extract_html_text(filepath: str) -> str:
    with open(filepath, "r", encoding="utf-8") as f:
        soup = BeautifulSoup(f.read(), "html.parser")
    for tag in soup(["script", "style"]):
        tag.decompose()
    return soup.get_text(separator="\n\n", strip=True)


def extract_txt_text(filepath: str) -> str:
    with open(filepath, "r", encoding="utf-8") as f:
        return f.read()


# Dispatch table: extension -> extractor function.
# Adding a new format later = one new function + one new entry here.
EXTRACTORS = {
    ".pdf": extract_pdf_text,
    ".docx": extract_docx_text,
    ".pptx": extract_pptx_text,
    ".html": extract_html_text,
    ".htm": extract_html_text,
    ".txt": extract_txt_text,
}


def extract_text(filepath: str) -> str | None:
    """Returns extracted text, or None if the file extension isn't supported."""
    for ext, extractor in EXTRACTORS.items():
        if filepath.lower().endswith(ext):
            return extractor(filepath)
    return None
